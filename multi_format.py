# multi_format.py
# Loaders for PDF / TXT / MD / VTT / SRT / DOCX / PPTX lecture content.
# Haofei Sun - CSE 5360

import re
from pathlib import Path


def load_text(path: str) -> str:
    with open(path, "r", encoding="utf-8", errors="ignore") as f:
        return f.read()


def load_pdf(path: str) -> str:
    try:
        from pypdf import PdfReader
    except ImportError:
        return load_text(path)
    reader = PdfReader(path)
    return "\n\n".join((page.extract_text() or "") for page in reader.pages)


def load_vtt(path: str) -> str:
    """WebVTT — drop timing lines and indices, keep speech."""
    raw = load_text(path)
    lines = []
    for line in raw.splitlines():
        line = line.strip()
        if "-->" in line:
            continue
        if line.upper().startswith("WEBVTT"):
            continue
        if re.match(r"^\d+$", line):
            continue
        if line:
            lines.append(line)
    return " ".join(lines)


def load_srt(path: str) -> str:
    """SRT — same as VTT but with index lines."""
    raw = load_text(path)
    lines = []
    for line in raw.splitlines():
        line = line.strip()
        if "-->" in line:
            continue
        if re.match(r"^\d+$", line):
            continue
        if line:
            lines.append(line)
    return " ".join(lines)


def load_docx(path: str) -> str:
    try:
        from docx import Document
    except ImportError:
        return ""
    doc = Document(path)
    return "\n".join(p.text for p in doc.paragraphs if p.text.strip())


def load_pptx(path: str) -> str:
    try:
        from pptx import Presentation
    except ImportError:
        return ""
    prs = Presentation(path)
    chunks = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                chunks.append(shape.text)
    return "\n".join(chunks)


def load_any(path: str) -> str:
    """Pick a loader by extension; fall back to plain text."""
    ext = Path(path).suffix.lower()
    if ext == ".pdf":
        return load_pdf(path)
    if ext in (".txt", ".md"):
        return load_text(path)
    if ext == ".vtt":
        return load_vtt(path)
    if ext == ".srt":
        return load_srt(path)
    if ext == ".docx":
        return load_docx(path)
    if ext == ".pptx":
        return load_pptx(path)
    return load_text(path)
